# -*- coding: utf-8 -*-
import json, re, os, collections
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

HERE = os.path.dirname(os.path.abspath(__file__))
DATA = os.path.join(HERE, 'ac_data.json')
OUT  = os.path.join(os.path.dirname(os.path.dirname(HERE)), '115冷氣機型比較_簡報.pptx')
d = json.load(open(DATA, encoding='utf-8'))
MAKETIME = '2026/07/10 09:12'
FONT = 'Microsoft JhengHei'

# ---- 色彩 (深色=高對比，黑白列印仍清楚) ----
TEAL   = RGBColor(0x0E,0x5C,0x63)   # 方式一主題深青
PLUM   = RGBColor(0x6D,0x2E,0x46)   # 方式二主題深莓
NAVY   = RGBColor(0x0B,0x4F,0x6C)   # 封面深藍
WHITE  = RGBColor(0xFF,0xFF,0xFF)
INK    = RGBColor(0x1A,0x1A,0x1A)   # 近黑內文
ZEBRA  = RGBColor(0xEE,0xF4,0xF5)
ZEBRA2 = RGBColor(0xF7,0xEF,0xF2)
PRICE_BG = RGBColor(0xFC,0xEF,0xD6) # 淡金 (列印成淺灰，文字仍深)
EFF_BG   = RGBColor(0xE4,0xF1,0xE4) # 淡綠
BORDER   = RGBColor(0x5A,0x5A,0x5A) # 深灰格線
FOOT     = RGBColor(0x66,0x66,0x66)

def num(v):
    if v in (None,''): return None
    if isinstance(v,(int,float)): return float(v)
    m=re.search(r'-?\d+(?:\.\d+)?',str(v).replace(',',''))
    return float(m.group()) if m else None
def money(v):
    n=num(v)
    return '—' if n is None else 'NT$ {:,}'.format(int(round(n)))

# ---- 功能標籤 ----
def feats(x, keep_heat=False):
    t=x['特色'] or ''; cat=x['細類'] or ''; out=[]
    if keep_heat:
        if '冷暖' in cat: out.append('冷暖')
        elif '冷專' in cat: out.append('冷專')
    if any(k in t for k in ['WiFi','Wi-Fi','wifi','智慧','APP','App','手機','語音','AI']): out.append('WiFi智慧')
    if any(k in t for k in ['自清','自動清','凍結洗淨','自體淨','霜效潔淨','洗淨']): out.append('自動清潔')
    if any(k in t for k in ['清淨','PM2.5','PM2','濾網','nanoe','Nanoe','等離子','光觸媒','過敏原']): out.append('空氣清淨')
    if any(k in t for k in ['抗菌','銀離子','防霉','抑菌']): out.append('抗菌防霉')
    if '除濕' in t: out.append('除濕')
    if any(k in t for k in ['防鏽','防腐','耐腐','耐蝕','防潮']): out.append('防鏽')
    if 'R32' in t or 'R-32' in t: out.append('R32冷媒')
    seen=set(); r=[]
    for g in out:
        if g not in seen: seen.add(g); r.append(g)
    return r
def feat_str(x, n=3):
    fs=feats(x)[:n]
    return '、'.join(fs) if fs else '—'

BRAND_SHORT={'SAMPO 聲寶':'聲寶','HERAN 禾聯':'禾聯','TECO 東元':'東元','Tatung 大同':'大同',
             'Panasonic':'Panasonic','三菱電機':'三菱電機','HITACHI 日立':'日立',
             'DAIKIN 大金':'大金','三菱重工':'三菱重工'}
def bshort(b): return BRAND_SHORT.get(b,b)
def type_short(cat):
    cat=cat or ''
    heat='冷暖' if '冷暖' in cat else ('冷專' if '冷專' in cat else '')
    if '窗型' in cat: b='窗型'
    elif '一對多' in cat: b='一對多'
    elif '吊隱' in cat: b='吊隱'
    elif '一對一' in cat: b='分離式'
    elif '移動' in cat: b='移動式'
    else: b=cat
    return b+('・'+heat if heat else '')

bt=collections.defaultdict(collections.Counter)
for x in d: bt[x['品牌']][x['細類']]+=1
brand_order=[b for b,_ in collections.Counter(x['品牌'] for x in d).most_common()]

# ================= PPTX 基礎 =================
prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW,SH=prs.slide_width,prs.slide_height

def add_slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(1,0,0,SW,SH); r.fill.solid(); r.fill.fore_color.rgb=bg
    r.line.fill.background(); r.shadow.inherit=False
    r._element.spPr.append(deepcopy(r._element.spPr.find(qn('a:noFill')))) if False else None
    # 移到最底
    sp=r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2,sp)
    return s

def txt(slide,l,t,w,h,text,size,color=INK,bold=False,align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,font=FONT,line=1.0,shrink=False):
    tb=slide.shapes.add_textbox(l,t,w,h); tf=tb.text_frame
    tf.word_wrap=True
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    tf.vertical_anchor=anchor
    lines=text.split('\n')
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=line
        r=p.add_run(); r.text=ln
        f=r.font; f.size=Pt(size); f.bold=bold; f.name=font; f.color.rgb=color
    return tb

def footer(slide, color=FOOT):
    txt(slide, Inches(0.35), Inches(7.06), Inches(12.63), Inches(0.34),
        '⚠ 僅供參考，商品規格/價格/贈品以中華電信系統實際受理與交貨為準，促銷價格會波動；圖片為示意，非實際外觀。   ｜   製表 '+MAKETIME,
        10.5, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---- 表格框線 ----
def set_border(cell,color=BORDER,w=9525):  # 9525 EMU = 0.75pt
    tcPr=cell._tc.get_or_add_tcPr()
    for tag in ('a:lnL','a:lnR','a:lnT','a:lnB'):
        e=tcPr.find(qn(tag))
        if e is not None: tcPr.remove(e)
        ln=tcPr.makeelement(qn(tag),{'w':str(w),'cap':'flat'})
        fill=ln.makeelement(qn('a:solidFill'),{})
        clr=ln.makeelement(qn('a:srgbClr'),{'val':'%02X%02X%02X'%(color[0],color[1],color[2])})
        fill.append(clr); ln.append(fill); tcPr.append(ln)

def cell_fmt(cell, text, size, color=INK, bold=False, align=PP_ALIGN.CENTER, bg=None):
    cell.margin_left=Pt(3); cell.margin_right=Pt(3); cell.margin_top=Pt(1); cell.margin_bottom=Pt(1)
    cell.vertical_anchor=MSO_ANCHOR.MIDDLE
    if bg is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb=bg
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb=WHITE
    tf=cell.text_frame; tf.word_wrap=True
    tf.clear()
    lines=str(text).split('\n')
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=0.95
        r=p.add_run(); r.text=ln
        f=r.font; f.size=Pt(size); f.bold=bold; f.name=FONT; f.color.rgb=color
    set_border(cell)

def make_table(slide, top, cols, widths, rows_data, header_bg, body_size=14, header_size=15,
               row_h=None, price_idx=None, eff_idx=None):
    ncol=len(cols); nrow=len(rows_data)+1
    total=sum(widths)
    left=Inches((13.333-total)/2.0); width=Inches(total)
    avail=Inches(7.0)-top
    if row_h is None:
        row_h=int(avail/ nrow)
    height=Emu(row_h*nrow)
    gf=slide.shapes.add_table(nrow,ncol,left,top,width,height)
    tbl=gf.table
    tbl.first_row=False; tbl.horz_banding=False
    # 關閉樣式帶狀
    for w,ci in zip(widths,range(ncol)):
        tbl.columns[ci].width=Inches(w)
    for ri in range(nrow):
        tbl.rows[ri].height=Emu(row_h)
    # header
    for ci,c in enumerate(cols):
        cell_fmt(tbl.cell(0,ci), c, header_size, color=WHITE, bold=True, bg=header_bg)
    # body
    for ri,row in enumerate(rows_data, start=1):
        zb = ZEBRA if header_bg==TEAL else ZEBRA2
        base = zb if (ri%2==0) else WHITE
        for ci,val in enumerate(row):
            bg=base; col=INK; bold=False; sz=body_size
            if price_idx is not None and ci==price_idx:
                bg=PRICE_BG; bold=True; col=INK; sz=body_size+1
            if eff_idx is not None and ci==eff_idx and str(val).strip() and '1級' in str(val):
                bg=EFF_BG; bold=True
            al=PP_ALIGN.LEFT if ci in (LEFTCOLS.get(id(cols),set())) else PP_ALIGN.CENTER
            cell_fmt(tbl.cell(ri,ci), val, sz, color=col, bold=bold, align=al, bg=bg)
    return gf

LEFTCOLS={}

def chunk(lst,n):
    return [lst[i:i+n] for i in range(0,len(lst),n)]

# ================= 封面 =================
def slide_cover():
    s=add_slide(NAVY)
    txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.5),
        '中華電信寬頻優惠專案', 22, color=RGBColor(0xCF,0xE6,0xEE), bold=True)
    txt(s, Inches(0.9), Inches(1.25), Inches(11.5), Inches(2.2),
        '115 年度智慧節能家電\n空調專區 · 冷氣機型比較表', 50, color=WHITE, bold=True, line=1.1)
    # 統計方塊
    stats=[('313','款上架機型'),('9','大品牌'),('2–10','KW 全涵蓋'),('2','種比較方式')]
    bx=Inches(0.9); bw=Inches(2.75); gap=Inches(0.18); by=Inches(3.7)
    for i,(nm,lb) in enumerate(stats):
        x=Emu(int(bx)+i*(int(bw)+int(gap)))
        card=s.shapes.add_shape(1,x,by,bw,Inches(1.5))
        card.fill.solid(); card.fill.fore_color.rgb=RGBColor(0x12,0x63,0x84)
        card.line.color.rgb=RGBColor(0x3E,0x8C,0xA8); card.line.width=Pt(1); card.shadow.inherit=False
        txt(s, x, Inches(3.85), bw, Inches(0.9), nm, 40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x, Inches(4.72), bw, Inches(0.4), lb, 15, color=RGBColor(0xCF,0xE6,0xEE), bold=True, align=PP_ALIGN.CENTER)
    # 警語框
    wb=s.shapes.add_shape(1, Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.95))
    wb.fill.solid(); wb.fill.fore_color.rgb=WHITE; wb.line.color.rgb=WHITE; wb.shadow.inherit=False
    txt(s, Inches(1.1), Inches(5.62), Inches(11.1), Inches(0.85),
        '⚠ 本表僅供參考：商品規格、價格與贈品以中華電信系統實際受理與交貨為準；促銷價格會波動。\n圖片為示意圖，非實際外觀，僅供辨識機型類別參考。',
        14.5, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE, line=1.15)
    txt(s, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.4),
        '收錄範圍：空調專區「上架中」機型　｜　製表時間 '+MAKETIME+'　｜　資料來源：（公告版）115年度智慧節能家電資料總表',
        11.5, color=RGBColor(0xB8,0xD4,0xDE))

# ================= 說明 / 圖例 =================
def slide_legend():
    s=add_slide(WHITE)
    txt(s, Inches(0.35), Inches(0.28), Inches(12.6), Inches(0.7),
        '如何使用本比較表', 34, color=NAVY, bold=True)
    # 左：兩種比較方式
    txt(s, Inches(0.5), Inches(1.15), Inches(6.0), Inches(0.5), '兩種比較方式', 22, color=TEAL, bold=True)
    m1=s.shapes.add_shape(1, Inches(0.5), Inches(1.75), Inches(5.9), Inches(1.55))
    m1.fill.solid(); m1.fill.fore_color.rgb=RGBColor(0xE9,0xF3,0xF4); m1.line.color.rgb=TEAL; m1.line.width=Pt(1.5); m1.shadow.inherit=False
    txt(s, Inches(0.7), Inches(1.85), Inches(5.5), Inches(1.4),
        '① 同噸數（KW／坪數）跨品牌比較\n　 先決定坪數，橫向比各品牌功能與價格。', 16, color=INK, bold=False, line=1.2, anchor=MSO_ANCHOR.MIDDLE)
    m2=s.shapes.add_shape(1, Inches(0.5), Inches(3.45), Inches(5.9), Inches(1.55))
    m2.fill.solid(); m2.fill.fore_color.rgb=RGBColor(0xF3,0xE9,0xEE); m2.line.color.rgb=PLUM; m2.line.width=Pt(1.5); m2.shadow.inherit=False
    txt(s, Inches(0.7), Inches(3.55), Inches(5.5), Inches(1.4),
        '② 同一廠商 2–10 KW 全系列\n　 選定品牌，縱向依坪數挑機型。', 16, color=INK, bold=False, line=1.2, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(0.5), Inches(5.2), Inches(5.9), Inches(1.6),
        '● 深色表頭、粗體價格，黑白列印也清楚。\n● 每頁皆有警語，以實際交貨為準。\n● 線上互動篩選版另見 HTML／GitHub Pages。',
        14.5, color=INK, line=1.3)
    # 右：功能標籤說明
    txt(s, Inches(6.9), Inches(1.15), Inches(6.0), Inches(0.5), '功能標籤說明', 22, color=NAVY, bold=True)
    tags=[('冷暖 / 冷專','是否具暖氣功能'),('R32冷媒','R32 新式環保冷媒'),
          ('WiFi智慧','手機App／語音操控'),('自動清潔','自動清洗／防霉'),
          ('空氣清淨','清淨濾網／PM2.5'),('抗菌防霉','抗菌銀離子／防霉'),
          ('除濕','可單獨除濕'),('防鏽','機體防鏽耐蝕')]
    gf=s.shapes.add_table(len(tags)+1,2, Inches(6.9), Inches(1.75), Inches(6.0), Inches(4.9))
    t=gf.table; t.columns[0].width=Inches(2.1); t.columns[1].width=Inches(3.9)
    t.first_row=False; t.horz_banding=False
    cell_fmt(t.cell(0,0),'標籤',15,color=WHITE,bold=True,bg=NAVY)
    cell_fmt(t.cell(0,1),'說明',15,color=WHITE,bold=True,bg=NAVY)
    for i,(a,b) in enumerate(tags,start=1):
        bg=ZEBRA if i%2==0 else WHITE
        cell_fmt(t.cell(i,0),a,15,color=INK,bold=True,bg=bg)
        cell_fmt(t.cell(i,1),b,14,color=INK,align=PP_ALIGN.LEFT,bg=bg)
    footer(s)

# ================= 分隔頁 =================
def slide_divider(no, title, sub, bg):
    s=add_slide(bg)
    txt(s, Inches(0.9), Inches(2.4), Inches(3.2), Inches(1.6), no, 120, color=RGBColor(0xFF,0xFF,0xFF), bold=True)
    txt(s, Inches(3.6), Inches(2.75), Inches(9.0), Inches(1.2), title, 46, color=WHITE, bold=True)
    txt(s, Inches(3.65), Inches(4.05), Inches(8.8), Inches(1.2), sub, 20, color=RGBColor(0xE7,0xEF,0xF2), line=1.25)

# ================= 方式一 =================
def build_method1():
    slide_divider('1', '比較方式一　同噸數跨品牌比較',
                  '依冷氣能力 KW（對應適用坪數）分級，將各品牌同級距機型並列，\n橫向比較功能與價格。', TEAL)
    cols=['品牌','型號','機型','冷氣\nkW','能效','CSPF','主要功能','市價','寬頻優惠價']
    widths=[1.15,2.30,1.45,0.60,0.90,0.62,2.40,1.30,1.40]  # 合計 12.12，置中
    LEFTCOLS[id(cols)]={1,6}
    BANDS=[('2.0–2.5 kW（約 3–5 坪）',1.9,2.59),('2.6–3.5 kW（約 5–7 坪）',2.6,3.59),
           ('3.6–4.5 kW（約 7–9 坪）',3.6,4.59),('4.6–5.9 kW（約 9–12 坪）',4.6,5.99),
           ('6.0–7.1 kW（約 12–15 坪）',6.0,7.19),('7.2–9.0 kW（約 15–20 坪）',7.2,9.09),
           ('9.1 kW 以上（約 20 坪以上）',9.1,9999)]
    PER=12
    for title,lo,hi in BANDS:
        items=[x for x in d if num(x['冷氣kW']) is not None and lo<=num(x['冷氣kW'])<=hi]
        items.sort(key=lambda x:(brand_order.index(x['品牌']) if x['品牌'] in brand_order else 99, num(x['冷氣kW']) or 0, num(x['優惠價']) or 0))
        parts=chunk(items,PER)
        for pi,part in enumerate(parts):
            s=add_slide(WHITE)
            suffix='' if len(parts)==1 else '　(%d/%d)'%(pi+1,len(parts))
            txt(s, Inches(0.35), Inches(0.22), Inches(9.5), Inches(0.62), title+suffix, 26, color=TEAL, bold=True)
            txt(s, Inches(9.9), Inches(0.30), Inches(3.08), Inches(0.5), '比較方式一・同噸數', 13, color=FOOT, bold=True, align=PP_ALIGN.RIGHT)
            rows=[]
            for x in part:
                rows.append([bshort(x['品牌']), x['型號'], type_short(x['細類']),
                             (('%g'%num(x['冷氣kW'])) if num(x['冷氣kW']) is not None else '—'),
                             x['能效等級'] or '—', x['CSPF'] or '—', feat_str(x),
                             money(x['市價']), money(x['優惠價'])])
            make_table(s, Inches(0.95), cols, widths, rows, TEAL, body_size=13, header_size=14,
                       price_idx=8, eff_idx=4)
            footer(s)

# ================= 方式二 =================
def build_method2():
    slide_divider('2', '比較方式二　同一廠商全系列',
                  '以品牌為單位，列出該廠商 2–10 KW 各機型，\n依冷氣能力由小到大，選定品牌後挑坪數。', PLUM)
    cols=['型號','機型','冷氣\nkW','暖氣\nkW','能效','CSPF','主要功能','市價','寬頻優惠價']
    widths=[2.60,1.45,0.60,0.60,0.90,0.62,2.60,1.30,1.40]  # 合計 12.07，置中
    LEFTCOLS[id(cols)]={0,6}
    PER=12
    for b in brand_order:
        items=[x for x in d if x['品牌']==b]
        items.sort(key=lambda x:(num(x['冷氣kW']) or 0, num(x['優惠價']) or 0))
        parts=chunk(items,PER)
        for pi,part in enumerate(parts):
            s=add_slide(WHITE)
            suffix='' if len(parts)==1 else '　(%d/%d)'%(pi+1,len(parts))
            # 品牌 chip
            chip=s.shapes.add_shape(5, Inches(0.35), Inches(0.22), Inches(3.7), Inches(0.62))
            chip.fill.solid(); chip.fill.fore_color.rgb=PLUM; chip.line.fill.background(); chip.shadow.inherit=False
            txt(s, Inches(0.35), Inches(0.22), Inches(3.7), Inches(0.62), b, 24, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            txt(s, Inches(4.2), Inches(0.30), Inches(5.5), Inches(0.5), '主力機型：'+type_short(bt[b].most_common(1)[0][0])+suffix, 15, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
            txt(s, Inches(9.9), Inches(0.30), Inches(3.08), Inches(0.5), '比較方式二・同廠商', 13, color=FOOT, bold=True, align=PP_ALIGN.RIGHT)
            rows=[]
            for x in part:
                rows.append([x['型號'], type_short(x['細類']),
                             (('%g'%num(x['冷氣kW'])) if num(x['冷氣kW']) is not None else '—'),
                             (('%g'%num(x['暖氣kW'])) if num(x['暖氣kW']) is not None else '—'),
                             x['能效等級'] or '—', x['CSPF'] or '—', feat_str(x),
                             money(x['市價']), money(x['優惠價'])])
            make_table(s, Inches(0.95), cols, widths, rows, PLUM, body_size=13, header_size=14,
                       price_idx=8, eff_idx=4)
            footer(s)

slide_cover()
slide_legend()
build_method1()
build_method2()
prs.save(OUT)
print('OK saved', OUT, 'slides=', len(prs.slides._sldIdLst))
